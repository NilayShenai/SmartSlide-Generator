# enhanced_advanced_ppt_maker.py

from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.prompts import PromptTemplate, FewShotPromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain.output_parsers import PydanticOutputParser
from pydantic import BaseModel, Field
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import ColorFormat
import re
from dotenv import load_dotenv
import os
import platform
import sys
import requests
import tempfile
from PIL import Image
import json
import base64
import webbrowser
from io import BytesIO
from playwright.sync_api import sync_playwright

# For PDF conversion
try:
    if platform.system() == "Windows":
        import comtypes.client
        COMTYPES_AVAILABLE = True
    else:
        COMTYPES_AVAILABLE = False
except ImportError:
    COMTYPES_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

load_dotenv()

# --- ENHANCED MODELS FOR IMAGES/FLOWCHARTS ---

class SlideContent(BaseModel):
    title: str = Field(description="The title of the slide")
    bullets: List[str] = Field(description="List of bullet points for the slide")
    image_query: Optional[str] = Field(default=None, description="Image search query for this slide (optional)")

class PresentationContent(BaseModel):
    slides: List[SlideContent] = Field(description="List of slides in the presentation")

# Enhanced examples with images and flowcharts
EXAMPLES = [
    {
        "topic": "Machine Learning Process",
        "num_slides": "3",
        "tone": "professional",
        "audience": "technical professionals",
        "theme": "technology",
        "output": """Slide 1: What is Machine Learning?
Machine Learning is a subset of artificial intelligence that enables systems to learn from data
Applications span across industries including healthcare, finance, and autonomous vehicles
Key paradigms include supervised, unsupervised, and reinforcement learning
Growing field with massive career opportunities and business impact
image_query: artificial intelligence machine learning concept

Slide 2: Machine Learning Workflow
Data collection and preprocessing from various sources
Feature engineering and selection for optimal model performance
Model training using algorithms like regression, decision trees, or neural networks
Model evaluation using metrics like accuracy, precision, and recall
Deployment and monitoring in production environments
flowchart_description: flowchart TD
    A[Data Collection] --> B[Data Preprocessing]
    B --> C[Feature Engineering]
    C --> D[Model Training]
    D --> E[Model Evaluation]
    E --> F[Good Performance?]
    F -->|No| C
    F -->|Yes| G[Model Deployment]

Slide 3: Popular ML Algorithms
Supervised learning: Linear regression, decision trees, random forests
Unsupervised learning: K-means clustering, PCA, hierarchical clustering
Deep learning: Convolutional Neural Networks (CNNs), Recurrent Neural Networks (RNNs)
Reinforcement learning: Q-learning, policy gradients
image_query: neural network diagram algorithms
"""
    },
    {
        "topic": "Digital Marketing Strategy",
        "num_slides": "2",
        "tone": "professional",
        "audience": "business executives",
        "theme": "corporate",
        "output": """Slide 1: Digital Marketing Landscape
Social media marketing drives 70% of brand engagement in 2024
Search engine optimization remains crucial for organic visibility
Email marketing delivers highest ROI at $42 for every $1 spent
Video content generates 1200% more shares than text and images combined
image_query: digital marketing social media analytics

Slide 2: Implementation Strategy
Phase 1: Audit current digital presence and competitor analysis
Phase 2: Develop content strategy and social media calendar
Phase 3: Launch paid advertising campaigns across multiple platforms
Phase 4: Monitor, analyze, and optimize based on performance metrics
flowchart_description: flowchart LR
    A[Audit] --> B[Strategy Development]
    B --> C[Campaign Launch]
    C --> D[Monitor & Optimize]
    D --> B
"""
    }
]

# Theme configurations (keeping your existing themes)
THEMES = {
    "corporate": {
        "colors": {
        "primary": RGBColor(26, 35, 126),     # Indigo (strong, modern)
        "secondary": RGBColor(144, 164, 174), # Blue-gray
        "accent": RGBColor(255, 87, 34)       # Deep orange accent
    },
    "background": {
        "type": "gradient",
        "colors": [RGBColor(236, 239, 241), RGBColor(207, 216, 220)]  # Sleek silver tones
    },
        "fonts": {"title": "Segoe UI Semibold", "body": "Segoe UI"},
    },

    "technology": {
        "colors": {
        "primary": RGBColor(0, 230, 118),     # Neon green
        "secondary": RGBColor(3, 169, 244),   # Bright blue
        "accent": RGBColor(255, 64, 129)      # Electric pink
    },
    "background": {
        "type": "gradient",
        "colors": [RGBColor(0, 0, 0), RGBColor(48, 48, 48)]  # Jet black to charcoal
    },
        "fonts": {"title": "Roboto", "body": "Roboto"},
    },

    "creative": {
        "colors": {
        "primary": RGBColor(255, 111, 145),   # Bubblegum pink
        "secondary": RGBColor(255, 202, 58),  # Gen Z yellow
        "accent": RGBColor(138, 201, 38)      # Lime green
    },
    "background": {
        "type": "gradient",
        "colors": [RGBColor(255, 243, 246), RGBColor(255, 224, 229)]  # Peach-pink fade
    },
        "fonts": {"title": "Poppins", "body": "Poppins"},
    },

    "academic": {
        "colors": {
        "primary": RGBColor(38, 70, 83),      # Keep: Elegant slate
        "secondary": RGBColor(244, 162, 97),  # Light terracotta
        "accent": RGBColor(231, 111, 81)      # Coral
    },
    "background": {
        "type": "gradient",
        "colors": [RGBColor(255, 251, 245), RGBColor(255, 241, 232)]
    },
        "fonts": {"title": "Georgia", "body": "Georgia"},
    },

    "health": {
        "colors": {
        "primary": RGBColor(102, 187, 106),   # Lively green
        "secondary": RGBColor(56, 79, 73),       # Dark desaturated green
        "accent": RGBColor(255, 143, 0)       # Sunlight orange
    },
    "background": {
        "type": "gradient",
        "colors": [RGBColor(255, 255, 255), RGBColor(230, 255, 244)]  # Soft wellness feel
    },
        "fonts": {"title": "Helvetica Neue", "body": "Helvetica"},
    },

    "environment": {
        "colors": {
        "primary": RGBColor(34, 87, 59),         # Dark forest green 
        "secondary": RGBColor(44, 106, 84),      # Deep jade green 
        "accent": RGBColor(255, 245, 157)     # Lemon zest
    },
    "background": {
        "type": "gradient",
        "colors": [RGBColor(232, 245, 233), RGBColor(200, 230, 201)]
    },
        "fonts": {"title": "Lato", "body": "Lato"},
    }
}

TEXT_SIZES = {
    "small": {"title": 26, "body": 16, "subtitle": 20},
    "medium": {"title": 30, "body": 20, "subtitle": 24},
    "large": {"title": 34, "body": 24, "subtitle": 28},
}

def extract_text_from_llm_output(output):
    """Extract text content from various LLM output types"""
    if isinstance(output, str):
        return output
    elif hasattr(output, 'content'):
        return output.content
    elif hasattr(output, 'text'):
        return output.text
    elif hasattr(output, 'message'):
        if hasattr(output.message, 'content'):
            return output.message.content
    else:
        return str(output)

def create_enhanced_few_shot_prompt():
    """Create enhanced few-shot prompt for image-only slides"""
    example_prompt = PromptTemplate(
        input_variables=["topic", "num_slides", "tone", "audience", "theme", "output"],
        template="""Topic: {topic}
Number of slides: {num_slides}
Tone: {tone}
Target audience: {audience}
Theme: {theme}

Output:
{output}"""
    )
    few_shot_prompt = FewShotPromptTemplate(
        examples=EXAMPLES,
        example_prompt=example_prompt,
        prefix="""You are an expert presentation creator with visual design expertise. For each slide, decide if:
1. A relevant image would enhance the content
2. A flowchart or diagram would help explain processes, workflows, or relationships

FORMAT FOR EACH SLIDE:
Slide N: Title
Bullet point 1 (without symbols - PowerPoint will add bullets automatically)
Bullet point 2 (clear, concise content)
Bullet point 3 (comprehensive information)
image_query: [specific search terms or omit if none]
flowchart_description: [mermaid.js flowchart syntax or omit if none]

CONTENT GUIDELINES:
- Write bullet points as plain text without manual bullet symbols (•, -, *)
- Keep bullets concise but informative (100-150 characters ideal)
- Focus on key points that support the slide title
- Ensure content is visually organized and professionally structured

FLOWCHART GUIDELINES:
- Use for processes, workflows, decision trees, or step-by-step procedures
- Use Mermaid.js syntax: flowchart TD (top-down) or flowchart LR (left-right)
- Example: flowchart TD
    A(["Start"]) --> B["Process"]
    B --> C["Decision?"]
    C -->|Yes| D["Action"]
    C -->|No| E["Alternative"]
    D --> F(["End"])
    E --> F

EXAMPLES:
""",
        suffix="""Topic: {topic}
Number of slides: {num_slides}
Tone: {tone}
Target audience: {audience}
Theme: {theme}

Create a {num_slides}-slide presentation following the exact format shown. Make intelligent decisions about visual elements (images and flowcharts) that will enhance understanding and engagement.

""",
        input_variables=["topic", "num_slides", "tone", "audience", "theme"]
    )
    return few_shot_prompt

def fetch_pexels_image(query, api_key, width=1200, height=800):
    """Enhanced Pexels API integration with better error handling and image quality"""
    if not api_key:
        print("⚠️ PEXELS_API_KEY not found in environment variables")
        return None
        
    url = f"https://api.pexels.com/v1/search"
    headers = {"Authorization": api_key}
    params = {
        "query": query,
        "per_page": 5,  # Get multiple options
        "orientation": "landscape",
        "size": "large"
    }
    
    try:
        print(f"🖼️ Searching for image: '{query}'")
        resp = requests.get(url, headers=headers, params=params, timeout=15)
        
        if resp.status_code == 200:
            data = resp.json()
            if data.get('photos'):
                # Try to get the best quality image
                photo = data['photos'][0]
                img_url = photo['src'].get('large2x', photo['src'].get('large', photo['src']['medium']))
                
                print(f"📥 Downloading image from Pexels...")
                img_resp = requests.get(img_url, timeout=15)
                
                if img_resp.status_code == 200:
                    # Save to temporary file
                    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.jpg')
                    tmp.write(img_resp.content)
                    tmp.close()
                    
                    # Resize image if needed using PIL
                    try:
                        with Image.open(tmp.name) as img:
                            # Resize to reasonable dimensions for PowerPoint
                            img = img.convert('RGB')
                            img.thumbnail((width, height), Image.Resampling.LANCZOS)
                            img.save(tmp.name, 'JPEG', quality=85)
                        
                        print(f"✅ Image downloaded and optimized: {tmp.name}")
                        return tmp.name
                    except Exception as e:
                        print(f"⚠️ Image processing failed: {e}")
                        return tmp.name  # Return original if processing fails
                        
        print(f"❌ No images found for query: '{query}'")
        return None
        
    except Exception as e:
        print(f"❌ Pexels API error: {e}")
        return None

def validate_mermaid_syntax(description):
    """Validate and clean Mermaid syntax"""
    # Remove empty lines and strip whitespace
    lines = [line.strip() for line in description.split('\n') if line.strip()]
    
    # Check if it's valid Mermaid syntax
    if not lines:
        return None
    
    # Skip if contains invalid keywords
    invalid_keywords = ['omit', 'none', 'null', 'undefined']
    for line in lines:
        if any(keyword in line.lower() for keyword in invalid_keywords):
            return None
    
    # Ensure it starts with a valid Mermaid diagram type
    first_line = lines[0].lower()
    valid_starts = ['flowchart', 'graph', 'sequencediagram', 'gantt', 'pie', 'gitgraph']
    
    if not any(first_line.startswith(start) for start in valid_starts):
        # If it doesn't start with a diagram type, assume it's a simple flowchart
        if len(lines) == 1 and '->' in lines[0]:
            return f"flowchart TD\n    {lines[0]}"
        return None
    
    return '\n'.join(lines)

def generate_mermaid_flowchart(description, width=800, height=600):
    """Generate a flowchart using Playwright and Mermaid.js"""
    try:
        print(f"📊 Generating flowchart: {description[:50]}...")
        
        # Validate and clean the Mermaid syntax
        cleaned_description = validate_mermaid_syntax(description)
        if not cleaned_description:
            print(f"❌ Invalid Mermaid syntax, skipping flowchart generation")
            return None
        
        print(f"🔍 Using cleaned description:\n{cleaned_description}")
        
        # Create HTML with Mermaid.js
        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <script src="https://cdn.jsdelivr.net/npm/mermaid@10.6.1/dist/mermaid.min.js"></script>
            <style>
                body {{ 
                    margin: 0; 
                    padding: 20px; 
                    background: white; 
                    font-family: Arial, sans-serif;
                }}
                #diagram {{ 
                    width: 100%; 
                    height: auto; 
                    display: flex; 
                    justify-content: center; 
                    align-items: center;
                    min-height: 400px;
                }}
                .mermaid {{
                    display: flex;
                    justify-content: center;
                    align-items: center;
                }}
            </style>
        </head>
        <body>
            <div id="diagram" class="mermaid">
{cleaned_description}
            </div>
            <script>
                console.log('Initializing Mermaid...');
                mermaid.initialize({{
                    startOnLoad: true,
                    theme: 'default',
                    securityLevel: 'loose',
                    fontFamily: 'Arial',
                    fontSize: 16,
                    flowchart: {{
                        useMaxWidth: true,
                        htmlLabels: true,
                        curve: 'basis'
                    }},
                    themeVariables: {{
                        primaryColor: '#ffffff',
                        primaryTextColor: '#333333',
                        primaryBorderColor: '#cccccc',
                        lineColor: '#666666'
                    }}
                }});
                
                // Manual rendering with callback
                window.addEventListener('load', function() {{
                    console.log('Page loaded, rendering Mermaid...');
                    const element = document.querySelector('#diagram');
                    if (element) {{
                        console.log('Found diagram element');
                        window.mermaidReady = true;
                    }} else {{
                        console.error('Diagram element not found');
                    }}
                }});
                
                // Error handling
                window.addEventListener('error', function(e) {{
                    console.error('JavaScript Error:', e.error);
                    window.mermaidError = e.error.toString();
                }});
            </script>
        </body>
        </html>
        """
        
        # Save to temporary HTML file
        tmp_html = tempfile.NamedTemporaryFile(delete=False, suffix='.html', mode='w', encoding='utf-8')
        tmp_html.write(html_content)
        tmp_html.close()
        
        print(f"📝 Created HTML file: {tmp_html.name}")
        
        # Use Playwright to capture screenshot
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True, args=['--no-sandbox', '--disable-dev-shm-usage'])
            page = browser.new_page()
            
            # Set viewport size
            page.set_viewport_size({"width": width, "height": height})
            
            # Load the HTML file
            print("🌐 Loading HTML file...")
            page.goto(f'file://{tmp_html.name}', wait_until='networkidle', timeout=30000)
            
            # Wait for page to be ready and Mermaid to initialize
            print("⏳ Waiting for Mermaid to render...")
            page.wait_for_timeout(3000)  # Give Mermaid more time to process
            
            # Check if Mermaid finished rendering
            ready = page.evaluate('() => window.mermaidReady || false')
            error = page.evaluate('() => window.mermaidError || null')
            
            if error:
                print(f"🔍 JavaScript error: {error}")
            
            if not ready:
                print("⚠️ Mermaid not ready, checking page content...")
                
            # Look for SVG or any rendered content
            svg_present = page.evaluate('() => !!document.querySelector("svg")')
            diagram_content = page.evaluate('() => document.querySelector("#diagram").innerHTML')
            
            print(f"🔍 SVG present: {svg_present}")
            print(f"🔍 Diagram content length: {len(diagram_content)}")
            
            if not svg_present:
                # Try alternative approach - wait for any content in the diagram div
                try:
                    page.wait_for_function(
                        '() => document.querySelector("#diagram").innerHTML.includes("svg") || document.querySelector("#diagram").innerHTML.length > 100',
                        timeout=10000
                    )
                except Exception as e:
                    print(f"⚠️ Timeout waiting for diagram content: {e}")
                    # Get final content for debugging
                    final_content = page.content()
                    print(f"🔍 Final page contains: SVG={('svg' in final_content.lower())}, Error={('error' in final_content.lower())}")
                    browser.close()
                    os.unlink(tmp_html.name)
                    return None
            
            # Create temporary file for screenshot
            tmp_img = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
            tmp_img_path = tmp_img.name
            tmp_img.close()
            
            # Take screenshot of the entire page
            print("📸 Taking screenshot...")
            page.screenshot(path=tmp_img_path, full_page=True)
            browser.close()
        
        # Clean up HTML file
        os.unlink(tmp_html.name)
        
        # Verify the image was created and has content
        if os.path.exists(tmp_img_path) and os.path.getsize(tmp_img_path) > 0:
            print(f"✅ Flowchart generated: {tmp_img_path}")
            return tmp_img_path
        else:
            print("❌ Generated image is empty or missing")
            if os.path.exists(tmp_img_path):
                os.unlink(tmp_img_path)
            return None
        
    except Exception as e:
        print(f"❌ Flowchart generation failed: {e}")
        # Clean up any temporary files
        if 'tmp_html' in locals() and os.path.exists(tmp_html.name):
            os.unlink(tmp_html.name)
        if 'tmp_img_path' in locals() and os.path.exists(tmp_img_path):
            os.unlink(tmp_img_path)
        return None
        return None

def parse_slides_enhanced(slide_text):
    """Parse slides for title, bullets, image_query, and flowchart_description."""
    slides = re.split(r'(?=Slide \d+:)', slide_text.strip())
    slides_data = []
    
    for slide_block in slides:
        if not slide_block.strip():
            continue
            
        lines = [line.strip() for line in slide_block.split('\n') if line.strip()]
        if not lines:
            continue
            
        # Extract title
        title_line = lines[0]
        title = title_line.split(': ', 1)[1] if ': ' in title_line else title_line
        title = re.sub(r'^Slide \d+:\s*', '', title).strip()
        
        # Initialize variables
        bullets = []
        image_query = None
        flowchart_description = None
        
        # Process remaining lines
        i = 1
        while i < len(lines):
            line = lines[i].strip()
            
            if line.lower().startswith('image_query:'):
                query = line.split(':', 1)[1].strip()
                image_query = query if query.lower() != 'none' else None
            elif line.lower().startswith('flowchart_description:'):
                # Handle multi-line flowchart description
                flowchart_content = line.split(':', 1)[1].strip()
                
                # Skip if it's "omit", "none", or other invalid values
                if flowchart_content.lower() in ['omit', 'none', 'null', '']:
                    pass  # Skip this flowchart
                else:
                    flowchart_lines = [flowchart_content]
                    i += 1
                    # Continue reading lines that are part of the flowchart
                    while i < len(lines) and not lines[i].strip().startswith(('image_query:', 'flowchart_description:', 'slide ')):
                        if lines[i].strip().startswith('-') or lines[i].strip().startswith('•'):
                            break  # This is a bullet point, stop reading flowchart
                        flowchart_lines.append(lines[i].strip())
                        i += 1
                    i -= 1  # Adjust for the extra increment
                    
                    # Join and validate the flowchart description
                    full_flowchart = '\n'.join(flowchart_lines).strip()
                    if full_flowchart and not any(word in full_flowchart.lower() for word in ['omit', 'none', 'null']):
                        flowchart_description = full_flowchart
            elif line.startswith('-') or line.startswith('•'):
                # Handle old format with manual bullet symbols
                bullet = re.sub(r'^[-•]\s*', '', line).strip()
                if bullet:
                    bullets.append(bullet)
            elif not line.lower().startswith(('image_query:', 'flowchart_description:', 'slide ')) and line.strip():
                # Handle new format without manual symbols (plain text bullets)
                # Skip empty lines and lines that are clearly not bullet points
                if len(line.strip()) > 5:  # Minimum reasonable bullet length
                    bullets.append(line.strip())
            
            i += 1
        
        slide_data = {
            'title': title,
            'bullets': bullets,
            'image_query': image_query,
            'flowchart_description': flowchart_description
        }
        
        slides_data.append(slide_data)
    
    return slides_data

def apply_background_to_slide(slide, background_config):
    """Apply background styling to a slide"""
    try:
        background = slide.background
        fill = background.fill
        if background_config["type"] == "solid":
            fill.solid()
            fill.fore_color.rgb = background_config["color"]
        elif background_config["type"] == "gradient":
            fill.gradient()
            fill.gradient_stops[0].color.rgb = background_config["colors"][0]
            if len(fill.gradient_stops) > 1:
                fill.gradient_stops[1].color.rgb = background_config["colors"][1]
    except Exception as e:
        print(f"⚠️ Could not apply background: {e}")

def apply_theme_to_slide(slide, theme_config, text_size_config):
    """Apply theme styling to a slide"""
    try:
        # Apply title formatting
        if slide.shapes.title:
            title_frame = slide.shapes.title.text_frame
            for paragraph in title_frame.paragraphs:
                paragraph.font.name = theme_config["fonts"]["title"]
                paragraph.font.size = Pt(text_size_config["title"])
                paragraph.font.color.rgb = theme_config["colors"]["primary"]
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER
        
        # Apply content formatting
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = theme_config["fonts"]["body"]
                    paragraph.font.size = Pt(text_size_config["body"])
                    paragraph.font.color.rgb = theme_config["colors"]["secondary"]
                    paragraph.space_before = Pt(6)
                    paragraph.space_after = Pt(6)
    except Exception as e:
        print(f"⚠️ Could not apply theme formatting: {e}")

def create_enhanced_ppt(slides_data, filename="presentation", output_format="pptx", theme="corporate", text_size="medium", flowcharts=None):
    """Create PowerPoint with image and flowchart visual elements"""
    prs = Presentation()
    theme_config = THEMES.get(theme, THEMES["corporate"])
    text_size_config = TEXT_SIZES.get(text_size, TEXT_SIZES["medium"])
    pexels_api_key = os.getenv("PEXELS_API_KEY")
    
    print(f"\n🎨 Creating presentation with theme: {theme}")
    
    # Create title slide
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    
    if slides_data:
        title_slide.shapes.title.text = slides_data[0]['title']
        if slides_data[0]['bullets']:
            title_slide.placeholders[1].text = slides_data[0]['bullets'][0]
    
    apply_background_to_slide(title_slide, theme_config["background"])
    apply_theme_to_slide(title_slide, theme_config, text_size_config)
    
    # Process content slides
    for i, slide_data in enumerate(slides_data[1:] if len(slides_data) > 1 else slides_data, 1):
        print(f"\n📄 Creating slide {i}: {slide_data['title']}")
        
        # Determine layout based on whether we have visuals
        has_visual = slide_data.get('flowchart_description') or slide_data.get('image_query')
        
        if has_visual:
            # Use two-content layout for slides with visuals
            slide = prs.slides.add_slide(prs.slide_layouts[3])  # Two content layout
            
            # Set title
            slide.shapes.title.text = slide_data['title']
            
            # Add bullet content to left placeholder
            if slide_data['bullets']:
                # Limit bullet points to prevent overflow
                max_bullets = 6  # Fewer bullets when we have visuals
                bullets_to_show = slide_data['bullets'][:max_bullets]
                
                # If we truncated bullets, add an indicator
                if len(slide_data['bullets']) > max_bullets:
                    bullets_to_show.append("...")
                
                # Use the left content placeholder which has built-in bullet formatting
                left_placeholder = slide.placeholders[1]  # Left content placeholder
                text_frame = left_placeholder.text_frame
                text_frame.clear()  # Clear default text
                
                # Add bullets using the built-in bullet system
                for j, bullet in enumerate(bullets_to_show):
                    # Truncate very long bullet points
                    max_bullet_length = 120
                    if len(bullet) > max_bullet_length:
                        bullet = bullet[:max_bullet_length-3] + "..."
                    
                    if j == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = bullet
                    p.level = 0  # First level bullets
                    
                    # Format the paragraph with enhanced spacing and bullet style
                    font = p.font
                    font.name = theme_config["fonts"]["body"]
                    font.size = Pt(text_size_config["body"])
                    font.color.rgb = theme_config["colors"]["secondary"]
                    
                    # Add professional spacing between bullets
                    p.space_before = Pt(6)
                    p.space_after = Pt(6)
                    p.line_spacing = 1.2
            
            # Add visual content to right placeholder
            right_placeholder = slide.placeholders[2]  # Right content placeholder
            
            # Add flowchart if present (prioritize flowchart)
            visual_added = False
            if slide_data.get('flowchart_description'):
                flowchart_path = generate_mermaid_flowchart(slide_data['flowchart_description'], width=400, height=300)
                if flowchart_path:
                    try:
                        # Replace the right placeholder with the flowchart image
                        slide.shapes._spTree.remove(right_placeholder._element)
                        visual_left = Inches(5.2)
                        visual_top = Inches(1.5)
                        visual_width = Inches(4.5)
                        slide.shapes.add_picture(flowchart_path, visual_left, visual_top, width=visual_width)
                        print(f"✅ Added flowchart to slide {i}")
                        visual_added = True
                        
                        # Clean up temporary file
                        try:
                            os.unlink(flowchart_path)
                        except:
                            pass
                            
                    except Exception as e:
                        print(f"❌ Failed to add flowchart to slide {i}: {e}")
            
            # Add image if present and no flowchart was added
            if not visual_added and slide_data.get('image_query'):
                img_path = fetch_pexels_image(slide_data['image_query'], pexels_api_key, width=600, height=400)
                if img_path:
                    try:
                        # Replace the right placeholder with the image
                        slide.shapes._spTree.remove(right_placeholder._element)
                        visual_left = Inches(5.2)
                        visual_top = Inches(1.5)
                        visual_width = Inches(4.5)
                        slide.shapes.add_picture(img_path, visual_left, visual_top, width=visual_width)
                        print(f"✅ Added image to slide {i}")
                        
                        # Clean up temporary file
                        try:
                            os.unlink(img_path)
                        except:
                            pass
                            
                    except Exception as e:
                        print(f"❌ Failed to add image to slide {i}: {e}")
            
            # Apply background and theme styling to visual slides
            apply_background_to_slide(slide, theme_config["background"])
            apply_theme_to_slide(slide, theme_config, text_size_config)
        
        else:
            # Use content with caption layout for text-only slides
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
            
            # Set title
            slide.shapes.title.text = slide_data['title']
            
            # Add bullet content using the built-in content placeholder
            if slide_data['bullets']:
                # More bullets allowed for text-only slides
                max_bullets = 8
                bullets_to_show = slide_data['bullets'][:max_bullets]
                
                # If we truncated bullets, add an indicator
                if len(slide_data['bullets']) > max_bullets:
                    bullets_to_show.append("...")
                
                # Use the content placeholder which has built-in bullet formatting
                content_placeholder = slide.placeholders[1]  # Content placeholder
                text_frame = content_placeholder.text_frame
                text_frame.clear()  # Clear default text
                
                # Add bullets using the built-in bullet system
                for j, bullet in enumerate(bullets_to_show):
                    # Truncate very long bullet points
                    max_bullet_length = 180
                    if len(bullet) > max_bullet_length:
                        bullet = bullet[:max_bullet_length-3] + "..."
                    
                    if j == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = bullet
                    p.level = 0  # First level bullets
                    
                    # Format the paragraph with enhanced spacing and bullet style
                    font = p.font
                    font.name = theme_config["fonts"]["body"]
                    font.size = Pt(text_size_config["body"])
                    font.color.rgb = theme_config["colors"]["secondary"]
                    
                    # Add professional spacing between bullets
                    p.space_before = Pt(6)
                    p.space_after = Pt(6)
                    p.line_spacing = 1.2
            
            # Apply background and theme styling to text-only slides
            apply_background_to_slide(slide, theme_config["background"])
            apply_theme_to_slide(slide, theme_config, text_size_config)
    
    # Add flowcharts after all regular slides if requested
    if flowcharts:
        for i, (flowchart_title, flowchart_desc) in enumerate(flowcharts, 1):
            print(f"\n📊 Creating flowchart slide {i}: {flowchart_title}")
            
            # Create a blank slide for the flowchart
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout
            
            # Add title at the top
            title_left = Inches(0.5)
            title_top = Inches(0.3)
            title_width = Inches(9)
            title_height = Inches(0.8)
            title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_shape.text = flowchart_title
            
            # Format title
            title_frame = title_shape.text_frame
            title_frame.margin_left = 0
            title_frame.margin_right = 0
            title_frame.margin_top = 0
            title_frame.margin_bottom = 0
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.font.name = theme_config["fonts"]["title"]
            title_paragraph.font.size = Pt(text_size_config["title"])
            title_paragraph.font.color.rgb = theme_config["colors"]["primary"]
            title_paragraph.font.bold = True
            title_paragraph.alignment = PP_ALIGN.CENTER
            
            # Generate and add flowchart (centered, larger size for dedicated flowchart slides)
            flowchart_path = generate_mermaid_flowchart(flowchart_desc, width=800, height=500)
            if flowchart_path:
                try:
                    # Center the flowchart on the slide
                    left = Inches(1.5)
                    top = Inches(1.8)
                    width = Inches(7)
                    slide.shapes.add_picture(flowchart_path, left, top, width=width)
                    print(f"✅ Added flowchart to slide")
                    
                    # Clean up temporary file
                    try:
                        os.unlink(flowchart_path)
                    except:
                        pass
                        
                except Exception as e:
                    print(f"❌ Failed to add flowchart: {e}")
            
            # Apply styling
            apply_background_to_slide(slide, theme_config["background"])
    
    # Save presentation
    if output_format.lower() == "pptx":
        # Use /tmp directory for Heroku compatibility
        if os.environ.get('DYNO'):  # Check if running on Heroku
            final_filename = f"/tmp/{filename}.pptx"
        else:
            final_filename = f"{filename}.pptx"
        prs.save(final_filename)
        print(f"\n✅ Enhanced presentation saved as '{final_filename}'")
    else:
        print(f"\n⚠️ Only PPTX format supported for enhanced presentations")
        # Use /tmp directory for Heroku compatibility
        if os.environ.get('DYNO'):  # Check if running on Heroku
            final_filename = f"/tmp/{filename}.pptx"
        else:
            final_filename = f"{filename}.pptx"
        prs.save(final_filename)
        print(f"📝 Saved as '{final_filename}'")

def get_user_inputs_enhanced(llm=None):
    """Collect all user inputs, using LLM to predict slide count if document is provided."""
    print("=== Enhanced PowerPoint Generator with AI Visuals ===\n")
    doc_path = input("Enter document path (.txt or .docx) to use as content (leave blank to AI generate it): ").strip().strip('"')
    topic = ""
    doc_content = None
    num_slides = None
    if doc_path:
        # Read document content
        ext = os.path.splitext(doc_path)[1].lower()
        if ext == ".txt":
            with open(doc_path, 'r', encoding='utf-8') as f:
                doc_content = f.read()
        elif ext == ".docx":
            from docx import Document as DocxDocument
            doc = DocxDocument(doc_path)
            doc_content = "\n".join([para.text for para in doc.paragraphs])
        else:
            print(f"Unsupported document type: {ext}. Only .txt and .docx are supported.")
            exit(1)
        if llm:
            print("Predicting suitable number of slides for your document...")
            num_slides = predict_num_slides(doc_content, llm)
            print(f"\n🔢 AI predicted number of slides: {num_slides}\n")
        else:
            num_slides = 5
    else:
        topic = input("Enter your presentation topic: ").strip()
        if not topic:
            topic = "General Topic"
        while True:
            try:
                num_slides = int(input("Number of slides (3-20, default: 5): ") or "5")
                if 3 <= num_slides <= 20:
                    break
                else:
                    print("Please enter a number between 3 and 20.")
            except ValueError:
                print("Please enter a valid number.")
    
    # Ask about flowcharts
    flowcharts = []
    while True:
        add_flowchart = input("\nDo you want to add a flowchart slide? (y/n): ").strip().lower()
        if add_flowchart == 'y':
            flowchart_title = input("Enter title for flowchart slide: ").strip()
            flowchart_desc = input("Enter flowchart description (e.g., 'A -> B -> C'): ").strip()
            flowcharts.append((flowchart_title, flowchart_desc))
        else:
            break
    
    # Enhanced tone selection
    tones = ["professional", "friendly", "urgent", "academic", "casual", "inspiring"]
    print(f"\nAvailable tones: {', '.join(tones)}")
    tone = input("Select tone (default: professional): ").strip().lower()
    if tone not in tones:
        tone = "professional"
    audiences = ["general public", "technical professionals", "students", "executives", "policymakers", "researchers"]
    print(f"\nSuggested audiences: {', '.join(audiences)}")
    audience = input("Target audience (default: general public): ").strip()
    if not audience:
        audience = "general public"
    available_themes = ["corporate", "creative", "academic", "health", "technology", "environment"]
    print(f"\nAvailable themes: {', '.join(available_themes)}")
    theme = input("Select theme (default: corporate): ").strip().lower()
    if theme not in available_themes:
        theme = "corporate"
    available_text_sizes = ["small", "medium", "large"]
    print(f"\nAvailable text sizes: {', '.join(available_text_sizes)}")
    text_size = input("Select text size (default: medium): ").strip().lower()
    if text_size not in available_text_sizes:
        text_size = "medium"
    formats = ["pptx", "pdf"]
    print(f"\nSupported formats: {', '.join(formats)}")
    output_format = input("Output format (default: pptx): ").strip().lower()
    if output_format not in formats:
        output_format = "pptx"
    filename = input("Filename (without extension, default: presentation): ").strip()
    if not filename:
        filename = "presentation"
    return {
        "doc_path": doc_path,
        "doc_content": doc_content,
        "topic": topic,
        "num_slides": str(num_slides),
        "tone": tone,
        "audience": audience,
        "theme": theme,
        "text_size": text_size,
        "output_format": output_format,
        "filename": filename,
        "flowcharts": flowcharts
    }

def predict_num_slides(doc_content, llm, min_slides=3, max_slides=20):
    """Use the LLM to predict a suitable number of slides for the document content."""
    prompt = f"""Given the following document content, suggest an appropriate number of PowerPoint slides (between {min_slides} and {max_slides}) for a presentation that covers all key points without being too brief or too detailed. Return only the number.\n\nDocument Content:\n{doc_content}\n"""
    result = llm.invoke(prompt)
    import re
    match = re.search(r"(\d+)", str(result))
    if match:
        num = int(match.group(1))
        return max(min_slides, min(num, max_slides))
    return min_slides

def main():
    """Main execution function"""
    try:
        # Check for required environment variables
        if not os.getenv("GOOGLE_API_KEY"):
            print("❌ Error: GOOGLE_API_KEY not found in environment variables")
            print("Please set your Google AI API key in the .env file")
            return
        
        # Initialize LLM
        print("🤖 Initializing AI model...")
        llm = ChatGoogleGenerativeAI(
            model="gemini-1.5-pro",
            temperature=0.0
        )
        
        # Get user inputs
        config = get_user_inputs_enhanced(llm)
        
        # Choose prompt based on input type
        if config["doc_path"]:
            # Use document summarization prompt
            summarization_prompt = PromptTemplate(
                input_variables=["doc_content", "num_slides", "tone", "audience", "theme"],
                template="""Summarize the following document content into a {num_slides}-slide PowerPoint presentation.
Each slide should have a title and 3-5 bullet points.
For each slide, decide if a relevant image or flowchart would enhance the content.

Format for each slide:
Slide N: Title
Bullet point 1 (without symbols - PowerPoint will add bullets automatically)
Bullet point 2 (clear, concise content)
Bullet point 3 (comprehensive information)
image_query: [specific search terms or omit if none]
flowchart_description: [mermaid.js flowchart syntax or omit if none]

CONTENT GUIDELINES:
- Write bullet points as plain text without manual bullet symbols (•, -, *)
- Keep bullets concise but informative (100-150 characters ideal)
- Focus on key points that support the slide title
- Ensure content is visually organized and professionally structured

FLOWCHART GUIDELINES:
- Use for processes, workflows, decision trees, or step-by-step procedures
- Use Mermaid.js syntax: flowchart TD (top-down) or flowchart LR (left-right)
- Example: flowchart TD
    A[Start] --> B[Process]
    B --> C{{Decision?}}
    C -->|Yes| D[Action]
    C -->|No| E[Alternative]

Tone: {tone}
Audience: {audience}
Theme: {theme}

Document Content:
{doc_content}

Return ONLY plain text in this format, as shown above.
"""
            )
            chain = summarization_prompt | llm
            output = chain.invoke({
                "doc_content": config["doc_content"],
                "num_slides": config["num_slides"],
                "tone": config["tone"],
                "audience": config["audience"],
                "theme": config["theme"]
            })
            output = extract_text_from_llm_output(output)
        else:
            # Ask user for topic and generate presentation on that topic
            topic = config["topic"]
            print(f"\nGenerating a presentation on the topic: {topic}")
            prompt_template = PromptTemplate(
                input_variables=["topic", "num_slides", "tone", "audience", "theme"],
                template="""Create a {num_slides}-slide PowerPoint presentation on the topic: {topic}.
Each slide should have a title and 3-5 bullet points.
For each slide, decide if a relevant image or flowchart would enhance the content.

Format for each slide:
Slide N: Title
Bullet point 1 (without symbols - PowerPoint will add bullets automatically)
Bullet point 2 (clear, concise content)
Bullet point 3 (comprehensive information)
image_query: [specific search terms or omit if none]
flowchart_description: [mermaid.js flowchart syntax or omit if none]

CONTENT GUIDELINES:
- Write bullet points as plain text without manual bullet symbols (•, -, *)
- Keep bullets concise but informative (100-150 characters ideal)
- Focus on key points that support the slide title
- Ensure content is visually organized and professionally structured

FLOWCHART GUIDELINES:
- Use for processes, workflows, decision trees, or step-by-step procedures
- Use Mermaid.js syntax: flowchart TD (top-down) or flowchart LR (left-right)
- Example: flowchart TD
    A[Start] --> B[Process]
    B --> C{{Decision?}}
    C -->|Yes| D[Action]
    C -->|No| E[Alternative]

Tone: {tone}
Audience: {audience}
Theme: {theme}

Return ONLY plain text in this format, as shown above.
"""
            )
            chain = prompt_template | llm
            output = chain.invoke({
                "topic": config["topic"],
                "num_slides": config["num_slides"],
                "tone": config["tone"],
                "audience": config["audience"],
                "theme": config["theme"]
            })
            output = extract_text_from_llm_output(output)
        
        print("✅ AI content generation complete!")
        print(f"📝 Generated content preview:\n{output[:200]}...\n")
        
        # Parse slides with enhanced visual information
        print("🔍 Parsing slide content and visual recommendations...")
        slides_data = parse_slides_enhanced(output)
        if not slides_data:
            print("❌ Error: Could not parse slide content. Please try again.")
            return
        print(f"📊 Parsed {len(slides_data)} slides successfully")
        
        # Display what visuals will be added
        print("\n🎨 Visual Enhancement Plan:")
        for i, slide in enumerate(slides_data, 1):
            visual_info = []
            if slide.get('image_query'):
                visual_info.append(f"Image: '{slide['image_query']}'")
            if slide.get('flowchart_description'):
                flowchart_preview = slide['flowchart_description'][:50].replace('\n', ' ')
                visual_info.append(f"Flowchart: {flowchart_preview}...")
            if visual_info:
                print(f"   Slide {i}: {' + '.join(visual_info)}")
            else:
                print(f"   Slide {i}: Text only")
        
        # Create the enhanced PowerPoint
        print(f"\n🚀 Creating enhanced PowerPoint presentation...")
        create_enhanced_ppt(
            slides_data=slides_data,
            filename=config['filename'],
            output_format="pptx",
            theme=config['theme'],
            text_size=config['text_size'],
            flowcharts=config.get('flowcharts', [])
        )
        
        print(f"\n🎉 Success! Your enhanced presentation is ready!")
        print(f"📁 File: {config['filename']}.pptx")
        print(f"🎨 Theme: {config['theme']}")
        print(f"📊 Total slides: {len(slides_data) + len(config.get('flowcharts', []))}")
        
        # Count visual enhancements
        image_count = sum(1 for slide in slides_data if slide.get('image_query'))
        llm_flowchart_count = sum(1 for slide in slides_data if slide.get('flowchart_description'))
        manual_flowchart_count = len(config.get('flowcharts', []))
        total_flowchart_count = llm_flowchart_count + manual_flowchart_count
        
        if image_count > 0:
            print(f"🖼️ Images added: {image_count}")
        if total_flowchart_count > 0:
            print(f"📊 Flowcharts added: {total_flowchart_count}")
            if llm_flowchart_count > 0:
                print(f"   - AI-generated: {llm_flowchart_count}")
            if manual_flowchart_count > 0:
                print(f"   - Manual: {manual_flowchart_count}")
    
    except KeyboardInterrupt:
        print("\n⚠️ Process interrupted by user")
    except Exception as e:
        print(f"❌ An error occurred: {str(e)}")
        print("Please check your API keys and try again")

def ensure_dependencies():
    """Check and install required dependencies"""
    required_packages = [
        ('langchain-google-genai', 'langchain_google_genai'),
        ('python-pptx', 'pptx'),
        ('python-dotenv', 'dotenv'),
        ('requests', 'requests'),
        ('pillow', 'PIL'),
        ('pydantic', 'pydantic'),
        ('playwright', 'playwright')
    ]
    
    missing_packages = []
    for pip_name, import_name in required_packages:
        try:
            __import__(import_name)
        except ImportError:
            missing_packages.append(pip_name)
    
    if missing_packages:
        print("❌ Missing required packages:")
        for package in missing_packages:
            print(f"   - {package}")
        print(f"\nInstall them with: pip install {' '.join(missing_packages)}")
        return False
    
    # Check if Playwright browsers are installed
    try:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as p:
            pass
    except Exception as e:
        print("⚠️ Playwright browsers not installed. Installing...")
        import sys
        import subprocess
        subprocess.run([sys.executable, "-m", "playwright", "install"], check=True)
    
    return True

def setup_environment():
    """Setup environment and check API keys"""
    print("🔧 Checking environment setup...")
    
    # Check for .env file
    if not os.path.exists('.env'):
        print("⚠️ .env file not found. Creating template...")
        with open('.env', 'w') as f:
            f.write("# Add your API keys here\n")
            f.write("GOOGLE_API_KEY=your_google_ai_api_key_here\n")
            f.write("PEXELS_API_KEY=your_pexels_api_key_here\n")
        print("📝 Created .env template. Please add your API keys and restart.")
        return False
    
    # Load environment variables
    load_dotenv()
    
    # Check required API keys
    google_key = os.getenv("GOOGLE_API_KEY")
    pexels_key = os.getenv("PEXELS_API_KEY")
    
    if not google_key or google_key == "your_google_ai_api_key_here":
        print("❌ GOOGLE_API_KEY not properly set in .env file")
        print("Get your key from: https://makersuite.google.com/app/apikey")
        return False
    
    if not pexels_key or pexels_key == "your_pexels_api_key_here":
        print("⚠️ PEXELS_API_KEY not set - images will not be added")
        print("Get your key from: https://www.pexels.com/api/")
    else:
        print("✅ Pexels API key configured")
    
    return True

if __name__ == "__main__":
    print("=== 🚀 Enhanced PowerPoint Generator with AI Visuals ===\n")
    
    # Check dependencies
    if not ensure_dependencies():
        exit(1)
    
    # Setup environment
    if not setup_environment():
        exit(1)
    
    # Run main program
    main()