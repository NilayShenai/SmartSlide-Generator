import os
import sys
from pathlib import Path
import platform
import re
import traceback

current_dir = Path(__file__).parent
parent_dir = current_dir.parent
sys.path.insert(0, str(parent_dir))

def convert_pptx_to_pdf(pptx_path: str, pdf_path: str):
    try:
        print(f"Starting PDF conversion: {pptx_path} -> {pdf_path}")
        
        # Try LibreOffice first (cross-platform, no watermarks, free)
        try:
            success = convert_pptx_to_pdf_libreoffice_enhanced(pptx_path, pdf_path)
            if success:
                print("PDF conversion successful using LibreOffice")
                return True
        except Exception as e:
            print(f"LibreOffice method failed: {e}")
        
        # Fallback to ReportLab-based extraction
        try:
            success = convert_pptx_to_pdf_with_reportlab(pptx_path, pdf_path)
            if success:
                print("PDF conversion successful using ReportLab")
                return True
        except Exception as e:
            print(f"ReportLab method failed: {e}")
        
        # Final fallback to simple text extraction
        try:
            success = convert_pptx_to_pdf_simple(pptx_path, pdf_path)
            if success:
                print("PDF conversion successful using simple text extraction")
                return True
        except Exception as e:
            print(f"Simple extraction method failed: {e}")
        
        raise Exception("All PDF conversion methods failed")
        
    except Exception as e:
        print(f"PDF conversion error: {e}")
        raise

def convert_pptx_to_pdf_libreoffice_enhanced(pptx_path: str, pdf_path: str):
    """
    Enhanced LibreOffice conversion with cross-platform support and improved reliability.
    Works on Windows, macOS, and Linux without watermarks.
    """
    try:
        import subprocess
        import shutil
        import tempfile
        
        print(f"Converting PPTX to PDF using LibreOffice: {pptx_path} -> {pdf_path}")
        
        # Verify input file exists
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"Source PPTX file not found: {pptx_path}")
        
        # Cross-platform LibreOffice detection
        libreoffice_cmd = detect_libreoffice_command()
        if not libreoffice_cmd:
            raise Exception("LibreOffice not found on system")
        
        # Verify LibreOffice is working
        if not verify_libreoffice(libreoffice_cmd):
            raise Exception("LibreOffice verification failed")
        
        # Create output directory if it doesn't exist
        pdf_dir = os.path.dirname(pdf_path)
        if pdf_dir:
            os.makedirs(pdf_dir, exist_ok=True)
        
        # Get absolute paths for reliable conversion
        pptx_abs_path = os.path.abspath(pptx_path)
        pdf_abs_path = os.path.abspath(pdf_path)
        
        # Use temporary directory for conversion
        with tempfile.TemporaryDirectory() as temp_dir:
            # Ensure clean LibreOffice state
            cleanup_libreoffice_processes()
            
            # Create automated LibreOffice profile to prevent first-run dialogs
            profile_dir = os.path.join(temp_dir, 'libreoffice_profile')
            create_automated_libreoffice_profile(profile_dir)
            
            # Build LibreOffice command
            cmd = build_libreoffice_command(libreoffice_cmd, pptx_abs_path, temp_dir)
            
            print(f"Executing LibreOffice command: {' '.join(cmd)}")
            
            try:
                # Update environment for temporary profile
                libreoffice_env = get_libreoffice_env()
                libreoffice_env['UserInstallation'] = f'file:///{profile_dir.replace(os.sep, "/")}'
                
                # Execute conversion with maximum automation and timeout
                result = subprocess.run(
                    cmd,
                    text=True,
                    timeout=180,                    # 3 minutes timeout
                    cwd=temp_dir,
                    stdin=subprocess.DEVNULL,       # Completely block interactive input
                    stdout=subprocess.PIPE,         # Capture all output
                    stderr=subprocess.PIPE,         # Capture all errors
                    env=libreoffice_env,            # Use our automated environment
                    shell=False,                    # Don't use shell for security
                    creationflags=subprocess.CREATE_NO_WINDOW if platform.system() == "Windows" else 0  # Hide window on Windows
                )
                
                if result.stdout:
                    print(f"LibreOffice stdout: {result.stdout.strip()}")
                if result.stderr and result.stderr.strip():
                    print(f"LibreOffice stderr: {result.stderr.strip()}")
                
                if result.returncode == 0:
                    # Find the generated PDF
                    base_name = os.path.splitext(os.path.basename(pptx_abs_path))[0]
                    generated_pdf = os.path.join(temp_dir, f"{base_name}.pdf")
                    
                    if os.path.exists(generated_pdf):
                        # Move to final destination
                        shutil.move(generated_pdf, pdf_abs_path)
                        
                        # Verify successful conversion
                        if os.path.exists(pdf_abs_path) and os.path.getsize(pdf_abs_path) > 1000:
                            file_size = os.path.getsize(pdf_abs_path)
                            print(f"✅ PDF conversion successful! 📄 Output file: {pdf_abs_path} ({file_size} bytes)")
                            return True
                        else:
                            print(f"❌ Generated PDF is too small or corrupt")
                            return False
                    else:
                        print(f"❌ Generated PDF not found at: {generated_pdf}")
                        # List files in temp directory for debugging
                        try:
                            temp_files = os.listdir(temp_dir)
                            print(f"Files in temp directory: {temp_files}")
                        except Exception:
                            pass
                        return False
                else:
                    print(f"❌ LibreOffice conversion failed with return code: {result.returncode}")
                    return False
                    
            except subprocess.TimeoutExpired:
                print("❌ LibreOffice conversion timed out after 3 minutes")
                cleanup_libreoffice_processes()  # Clean up hanging processes
                return False
            except Exception as subprocess_error:
                print(f"❌ Subprocess error: {subprocess_error}")
                cleanup_libreoffice_processes()  # Clean up hanging processes
                return False
            finally:
                # Always attempt cleanup after conversion
                cleanup_libreoffice_processes()
        
    except Exception as e:
        print(f"❌ LibreOffice conversion failed: {e}")
        return False

def detect_libreoffice_command():
    """
    Detect LibreOffice command across different platforms.
    Enhanced for Linux VPS environments with multiple installation methods.
    """
    import shutil
    
    system = platform.system()
    
    # Platform-specific LibreOffice paths and commands
    if system == "Darwin":  # macOS
        possible_paths = [
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "/opt/homebrew/bin/soffice",
            "/usr/local/bin/soffice"
        ]
        for path in possible_paths:
            if os.path.exists(path):
                print(f"🔍 Found LibreOffice on macOS: {path}")
                return path
        
        # Try PATH lookup
        for cmd in ['soffice', 'libreoffice']:
            if shutil.which(cmd):
                print(f"🔍 Found LibreOffice in PATH: {cmd}")
                return cmd
                
    elif system == "Linux":
        # Enhanced Linux VPS paths (apt, snap, flatpak, manual installs)
        possible_paths = [
            "/usr/bin/libreoffice",
            "/usr/bin/soffice", 
            "/usr/local/bin/libreoffice",
            "/usr/local/bin/soffice",
            "/opt/libreoffice/program/soffice",
            "/snap/bin/libreoffice",  # Snap package
            "/var/lib/flatpak/exports/bin/org.libreoffice.LibreOffice",  # System flatpak
            os.path.expanduser("~/.local/share/flatpak/exports/bin/org.libreoffice.LibreOffice"),  # User flatpak
        ]
        
        for path in possible_paths:
            if os.path.exists(path):
                print(f"🔍 Found LibreOffice on Linux VPS: {path}")
                return path
        
        # Try PATH lookup with various command names (for different versions)
        for cmd in ['libreoffice', 'soffice', 'libreoffice7.0', 'libreoffice6.4', 'libreoffice7.2', 'libreoffice7.4']:
            if shutil.which(cmd):
                print(f"🔍 Found LibreOffice in PATH: {cmd}")
                return cmd
                
    elif system == "Windows":
        # Windows-specific paths
        possible_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            r"C:\Program Files\LibreOffice 7\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice 7\program\soffice.exe",
        ]
        for path in possible_paths:
            if os.path.exists(path):
                print(f"🔍 Found LibreOffice on Windows: {path}")
                return path
        
        # Try PATH lookup
        for cmd in ['soffice', 'soffice.exe', 'libreoffice', 'libreoffice.exe']:
            if shutil.which(cmd):
                print(f"🔍 Found LibreOffice in PATH: {cmd}")
                return cmd
    
    print("❌ LibreOffice not found on system")
    if system == "Linux":
        print("💡 For Linux VPS, install with: sudo apt update && sudo apt install libreoffice --no-install-recommends")
    return None

def verify_libreoffice(libreoffice_cmd):
    """
    Verify that LibreOffice is working properly.
    """
    try:
        import subprocess
        result = subprocess.run(
            [libreoffice_cmd, '--version'],
            capture_output=True,
            text=True,
            timeout=30
        )
        
        if result.returncode == 0:
            version_info = result.stdout.strip() if result.stdout else "LibreOffice (version check passed)"
            print(f"✅ LibreOffice verified: {version_info}")
            return True
        else:
            print(f"❌ LibreOffice verification failed: return code {result.returncode}")
            if result.stderr:
                print(f"Error output: {result.stderr.strip()}")
            return False
            
    except Exception as e:
        print(f"❌ LibreOffice verification error: {e}")
        return False

def build_libreoffice_command(libreoffice_cmd, input_file, output_dir):
    """
    Build LibreOffice command with proven automation parameters that work reliably.
    Optimized for headless Linux VPS and local environments.
    """
    cmd = [
        libreoffice_cmd,
        '--headless',                    # No GUI - ESSENTIAL
        '--invisible',                   # Hide splash screen
        '--nolockcheck',                 # Don't check for lock files
        '--nologo',                      # Don't show logo
        '--norestore',                   # Don't restore windows
        '--nofirststartwizard',          # Skip first start wizard
        '--convert-to', 'pdf',           # Convert to PDF format
        '--outdir', output_dir,          # Output directory
        input_file                       # Input file
    ]
    
    # Add Linux VPS-specific parameters for maximum compatibility
    import platform
    if platform.system() == "Linux":
        cmd.insert(1, '--nodefault')     # Don't open default document
        cmd.insert(2, '--nocrashreport') # Disable crash reporting
        cmd.insert(3, '--disable-extension-update') # Disable extension updates
    
    return cmd

def get_libreoffice_env():
    """
    Create a clean environment for LibreOffice with essential automation variables.
    Optimized for headless Linux VPS environments.
    """
    import os
    env = os.environ.copy()
    
    # Set essential environment variables for non-interactive operation
    env.update({
        'SAL_USE_VCLPLUGIN': 'svp',                  # Use Server Virtual Plugin (headless)
        'SAL_NO_MOUSEGRABS': '1',                    # Prevent mouse grabs
        'SAL_DISABLE_CRASHREPORT': '1',              # Disable crash reporting
    })
    
    # Platform-specific essential settings
    system = platform.system()
    if system == "Linux":
        # Enhanced Linux VPS support - headless server environment
        env.update({
            'DISPLAY': '',                           # No X11 display (headless)
            'HOME': '/tmp/libreoffice_vps_home',     # Temporary home for VPS
            'XDG_RUNTIME_DIR': '/tmp/xdg_runtime',   # XDG runtime directory
            'TMPDIR': '/tmp',                        # Temp directory
            'FONTCONFIG_PATH': '/etc/fonts',         # Font configuration
            'QT_QPA_PLATFORM': 'offscreen',          # Qt offscreen platform
            'GDK_BACKEND': 'x11',                    # GTK backend
        })
        
        # Create necessary directories for VPS
        try:
            os.makedirs('/tmp/libreoffice_vps_home', exist_ok=True)
            os.makedirs('/tmp/xdg_runtime', exist_ok=True)
        except:
            pass
            
    elif system == "Windows":
        env.update({
            'TEMP': env.get('TEMP', 'C:\\Temp'),
            'TMP': env.get('TMP', 'C:\\Temp'),
        })
    
    return env

def cleanup_libreoffice_processes():
    """
    Clean up any hanging LibreOffice processes to prevent conflicts.
    """
    try:
        import subprocess
        system = platform.system()
        
        if system == "Windows":
            # Kill LibreOffice processes on Windows
            subprocess.run([
                'taskkill', '/F', '/IM', 'soffice.exe'
            ], capture_output=True, timeout=10)
            subprocess.run([
                'taskkill', '/F', '/IM', 'soffice.bin'
            ], capture_output=True, timeout=10)
        else:
            # Kill LibreOffice processes on Linux/macOS
            subprocess.run([
                'pkill', '-f', 'soffice'
            ], capture_output=True, timeout=10)
            subprocess.run([
                'pkill', '-f', 'libreoffice'
            ], capture_output=True, timeout=10)
            
    except Exception as e:
        # Process cleanup is best-effort, don't fail if it doesn't work
        print(f"Process cleanup warning: {e}")
        pass

def create_automated_libreoffice_profile(profile_dir):
    """
    Create a pre-configured LibreOffice profile to prevent first-run dialogs.
    """
    try:
        import os
        
        # Create profile directory structure
        os.makedirs(profile_dir, exist_ok=True)
        os.makedirs(os.path.join(profile_dir, 'user'), exist_ok=True)
        os.makedirs(os.path.join(profile_dir, 'user', 'config'), exist_ok=True)
        
        # Create registrymodifications.xcu to skip first-run wizard
        registry_config = '''<?xml version="1.0" encoding="UTF-8"?>
<oor:items xmlns:oor="http://openoffice.org/2001/registry" xmlns:xs="http://www.w3.org/2001/XMLSchema" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
<item oor:path="/org.openoffice.Setup/Office">
<prop oor:name="FirstStartWizardCompleted" oor:op="fuse">
<value>true</value>
</prop>
</item>
<item oor:path="/org.openoffice.Office.Common/Misc">
<prop oor:name="FirstRun" oor:op="fuse">
<value>false</value>
</prop>
</item>
</oor:items>'''
        
        registry_path = os.path.join(profile_dir, 'user', 'registrymodifications.xcu')
        with open(registry_path, 'w', encoding='utf-8') as f:
            f.write(registry_config)
            
        print(f"✅ Created automated LibreOffice profile at: {profile_dir}")
        return True
        
    except Exception as e:
        print(f"⚠️ Warning: Could not create LibreOffice profile: {e}")
        return False

def convert_pptx_to_pdf_simple(pptx_path: str, pdf_path: str):
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import inch
        from pptx import Presentation
        
        print(f"Converting PPTX to PDF using simple text extraction: {pptx_path} -> {pdf_path}")
        
        prs = Presentation(pptx_path)
        
        pdf_dir = os.path.dirname(pdf_path)
        if pdf_dir:  # Only create directory if it's not empty
            os.makedirs(pdf_dir, exist_ok=True)
        
        c = canvas.Canvas(pdf_path, pagesize=A4)
        page_width, page_height = A4
        
        for slide_idx, slide in enumerate(prs.slides):
            if slide_idx > 0:
                c.showPage()
            
            c.setFont("Helvetica-Bold", 16)
            c.drawString(inch, page_height - inch, f"Slide {slide_idx + 1}")
            
            y_position = page_height - 1.5 * inch
            c.setFont("Helvetica", 11)
            
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text = shape.text.strip()
                        lines = text.split('\n')
                        
                        for line in lines:
                            if line.strip() and y_position > inch:
                                if len(line) > 80:
                                    words = line.split(' ')
                                    current_line = ''
                                    for word in words:
                                        if len(current_line + ' ' + word) < 80:
                                            current_line += (' ' if current_line else '') + word
                                        else:
                                            if current_line:
                                                c.drawString(inch, y_position, current_line)
                                                y_position -= 15
                                            current_line = word
                                    if current_line and y_position > inch:
                                        c.drawString(inch, y_position, current_line)
                                        y_position -= 15
                                else:
                                    c.drawString(inch, y_position, line.strip())
                                    y_position -= 15
                
                except Exception as shape_error:
                    print(f"Error processing shape: {shape_error}")
                    continue
            
            c.setFont("Helvetica", 8)
            c.drawString(page_width - 1.5*inch, 0.5*inch, f"Page {slide_idx + 1}")
        
        c.save()
        
        if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 500:
            print(f"Successfully converted to PDF using simple method: {pdf_path}")
            return True
        else:
            raise Exception("PDF file was not created properly")
            
    except Exception as e:
        print(f"Simple PDF conversion failed: {e}")
        return False

def convert_pptx_to_pdf_with_reportlab(pptx_path: str, pdf_path: str):
    try:
        from pptx import Presentation
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.lib.utils import ImageReader
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib.colors import black, blue
        import io
        from PIL import Image
        
        print(f"Converting PPTX to PDF using ReportLab: {pptx_path} -> {pdf_path}")
        
        prs = Presentation(pptx_path)
        
        doc = SimpleDocTemplate(pdf_path, pagesize=A4)
        story = []
        styles = getSampleStyleSheet()
        
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=12,
            textColor=blue
        )
        
        bullet_style = ParagraphStyle(
            'CustomBullet',
            parent=styles['Normal'],
            fontSize=12,
            spaceAfter=6,
            leftIndent=20,
            bulletIndent=10
        )
        
        for slide_num, slide in enumerate(prs.slides, 1):
            story.append(Paragraph(f"Slide {slide_num}", title_style))
            story.append(Spacer(1, 12))
            
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            
            for text in slide_text:
                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        story.append(Paragraph(line.strip(), bullet_style))
            
            if slide_num < len(prs.slides):
                story.append(PageBreak())
        
        doc.build(story)
        
        if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 1000:
            print(f"Successfully converted to PDF using ReportLab: {pdf_path}")
            return True
        else:
            raise Exception("PDF file was not created properly")
            
    except ImportError as e:
        print(f"ReportLab not available: {e}")
        return False
    except Exception as e:
        print(f"ReportLab conversion failed: {e}")
        return False

def convert_pptx_to_pdf_com(pptx_path: str, pdf_path: str):
    try:
        import comtypes.client
        import comtypes
        import time
        
        comtypes.CoInitialize()
        
        pptx_abs_path = os.path.abspath(pptx_path)
        pdf_abs_path = os.path.abspath(pdf_path)
        
        if not os.path.exists(pptx_abs_path):
            raise Exception(f"Source PPTX file not found: {pptx_abs_path}")
        
        print(f"Converting {pptx_abs_path} to {pdf_abs_path}")
        
        powerpoint = None
        presentation = None
        
        try:
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = False
            powerpoint.DisplayAlerts = False
            time.sleep(1)
            
            presentation = powerpoint.Presentations.Open(
                pptx_abs_path, 
                ReadOnly=True, 
                Untitled=False, 
                WithWindow=False
            )
            time.sleep(1)
            
            pdf_dir = os.path.dirname(pdf_abs_path)
            os.makedirs(pdf_dir, exist_ok=True)
            
            presentation.ExportAsFixedFormat(
                OutputFileName=pdf_abs_path,
                FixedFormatType=2,
                Intent=1,
                FrameSlides=False,
                HandoutOrder=1,
                OutputType=1,
                PrintHiddenSlides=False,
                PrintRange=None,
                RangeType=1,
                SlideShowName="",
                IncludeDocProps=True,
                KeepIRMSettings=True,
                DocStructureTags=True,
                BitmapMissingFonts=True,
                UseDocumentICCProfile=False
            )
            
            time.sleep(2)
            
            if os.path.exists(pdf_abs_path) and os.path.getsize(pdf_abs_path) > 1000:
                print(f"Successfully converted to PDF using COM: {pdf_abs_path}")
                return True
            else:
                raise Exception("PDF file was not created properly")
                
        finally:
            try:
                if presentation:
                    presentation.Close()
                if powerpoint:
                    powerpoint.Quit()
                time.sleep(1)
            except:
                pass
            
            try:
                comtypes.CoUninitialize()
            except:
                pass
                
    except ImportError:
        print("comtypes not available for PDF conversion")
        return False
    except Exception as e:
        print(f"COM automation failed: {e}")
        return False

def convert_pptx_to_pdf_libreoffice(pptx_path: str, pdf_path: str):
    try:
        import subprocess
        import shutil
        
        libreoffice_cmd = None
        for cmd in ['libreoffice', 'soffice']:
            if shutil.which(cmd):
                libreoffice_cmd = cmd
                break
        
        if not libreoffice_cmd:
            print("LibreOffice not found in PATH")
            return False
            
        try:
            result = subprocess.run([libreoffice_cmd, '--version'], 
                                    capture_output=True, check=True, timeout=10)
            print(f"Found LibreOffice: {result.stdout.decode().strip()}")
        except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
            print("LibreOffice not available or not responding")
            return False
        
        print(f"Converting PPTX to PDF using LibreOffice: {pptx_path} -> {pdf_path}")
        
        pptx_abs_path = os.path.abspath(pptx_path)
        pdf_abs_path = os.path.abspath(pdf_path)
        pdf_dir = os.path.dirname(pdf_abs_path)
        
        os.makedirs(pdf_dir, exist_ok=True)
        
        import tempfile
        with tempfile.TemporaryDirectory() as temp_dir:
            cmd = [
                libreoffice_cmd,
                '--headless',
                '--invisible',
                '--nodefault', 
                '--nolockcheck',
                '--nologo',
                '--norestore',
                '--convert-to', 'pdf',
                '--outdir', temp_dir,
                pptx_abs_path
            ]
            
            try:
                result = subprocess.run(
                    cmd, 
                    capture_output=True, 
                    text=True, 
                    timeout=120,
                    cwd=temp_dir
                )
                
                print(f"LibreOffice command: {' '.join(cmd)}")
                print(f"Return code: {result.returncode}")
                if result.stdout:
                    print(f"Stdout: {result.stdout}")
                if result.stderr:
                    print(f"Stderr: {result.stderr}")
                
                if result.returncode == 0:
                    base_name = os.path.splitext(os.path.basename(pptx_abs_path))[0]
                    generated_pdf = os.path.join(temp_dir, f"{base_name}.pdf")
                    
                    print(f"Looking for generated PDF: {generated_pdf}")
                    
                    if os.path.exists(generated_pdf):
                        shutil.move(generated_pdf, pdf_abs_path)
                        
                        if os.path.exists(pdf_abs_path) and os.path.getsize(pdf_abs_path) > 1000:
                            print(f"Successfully converted to PDF using LibreOffice: {pdf_abs_path}")
                            return True
                        else:
                            print(f"Generated PDF is too small: {os.path.getsize(pdf_abs_path) if os.path.exists(pdf_abs_path) else 'not found'}")
                    else:
                        print(f"Generated PDF not found at expected location: {generated_pdf}")
                        try:
                            temp_files = os.listdir(temp_dir)
                            print(f"Files in temp directory: {temp_files}")
                        except:
                            pass
                else:
                    print(f"LibreOffice conversion failed with return code: {result.returncode}")
            
            except subprocess.TimeoutExpired:
                print("LibreOffice conversion timed out")
            except Exception as subprocess_error:
                print(f"Subprocess error: {subprocess_error}")
        
        return False
        
    except Exception as e:
        print(f"LibreOffice conversion failed: {e}")
        return False

try:
    from dotenv import load_dotenv
    load_dotenv(parent_dir / '.env')
except ImportError:
    print("Warning: python-dotenv not installed, environment variables might not load")

import json
import tempfile
import threading
import uuid
from datetime import datetime, timedelta
import logging
from typing import Dict, List, Optional, Any
import traceback

try:
    from ppt_maker_flo import (
        ChatGoogleGenerativeAI,
        create_enhanced_ppt,
        parse_slides_enhanced,
        extract_text_from_llm_output,
        create_enhanced_few_shot_prompt,
        predict_num_slides,
        THEMES,
        TEXT_SIZES
    )
except ImportError as e:
    print(f"Error importing ppt_maker_flo: {e}")
    print("Current working directory:", os.getcwd())
    print("Python path:", sys.path)
    raise

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

if os.environ.get('DYNO'):
    import logging
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s %(levelname)s: %(message)s [in %(pathname)s:%(lineno)d]'
    )
    
    import psutil
    logger.info(f"Memory usage: {psutil.virtual_memory().percent}%")
    logger.info(f"Available memory: {psutil.virtual_memory().available / 1024 / 1024:.1f} MB")

active_jobs: Dict[str, Dict] = {}
job_results: Dict[str, Dict] = {}

ALLOWED_EXTENSIONS = {'txt', 'docx', 'pdf'}

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def cleanup_old_files():
    try:
        current_time = datetime.now()
        cutoff_time = current_time - timedelta(hours=24)
        
        if os.environ.get('DYNO'):
            upload_folder = '/tmp'
            output_folder = '/tmp'
        else:
            upload_folder = 'uploads'
            output_folder = 'outputs'
        
        if os.path.exists(upload_folder):
            for filename in os.listdir(upload_folder):
                filepath = os.path.join(upload_folder, filename)
                if os.path.isfile(filepath):
                    file_time = datetime.fromtimestamp(os.path.getmtime(filepath))
                    if file_time < cutoff_time:
                        os.remove(filepath)
                        logger.info(f"Cleaned up old upload: {filename}")
        
        if os.path.exists(output_folder):
            for filename in os.listdir(output_folder):
                filepath = os.path.join(output_folder, filename)
                if os.path.isfile(filepath):
                    file_time = datetime.fromtimestamp(os.path.getmtime(filepath))
                    if file_time < cutoff_time:
                        os.remove(filepath)
                        logger.info(f"Cleaned up old output: {filename}")
                        
    except Exception as e:
        logger.error(f"Error during cleanup: {e}")

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF using pdfplumber for better accuracy."""
    try:
        import pdfplumber
        text = ""
        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        
        final_text = text.strip()
        if not final_text:
            logger.warning(f"No extractable text found in PDF: {file_path}. This may be an image-based PDF.")
            raise Exception("This PDF appears to contain only images without extractable text. Please try a PDF with text content or convert it to a text-based format.")
        
        return final_text
    except ImportError:
        logger.error("pdfplumber not installed for PDF support")
        raise ImportError("pdfplumber not installed. Please install with: pip install pdfplumber")
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        raise Exception(f"Error reading PDF file: {str(e)}")

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX using python-docx."""
    try:
        from docx import Document
        doc = Document(file_path)
        text_parts = []
        
        # Extract text from paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_parts.append(cell.text.strip())
        
        return "\n".join(text_parts)
    except ImportError:
        logger.error("python-docx not installed for DOCX support")
        raise ImportError("python-docx not installed. Please install with: pip install python-docx")
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {e}")
        raise Exception(f"Error reading DOCX file: {str(e)}")

def extract_text_from_txt(file_path: str) -> str:
    """Extract text from plain text file."""
    try:
        # Try multiple encodings
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    return f.read().strip()
            except UnicodeDecodeError:
                continue
        
        # If all encodings fail, try with error handling
        with open(file_path, "r", encoding='utf-8', errors='replace') as f:
            return f.read().strip()
            
    except Exception as e:
        logger.error(f"Error reading text file: {e}")
        raise Exception(f"Error reading text file: {str(e)}")

def detect_file_type(file_path: str) -> str:
    """Detect file type using multiple methods."""
    try:
        # First try python-magic if available
        try:
            import magic
            mime = magic.Magic(mime=True)
            return mime.from_file(file_path)
        except ImportError:
            pass
        except Exception as e:
            logger.warning(f"python-magic detection failed: {e}")
    
        # Fallback to file extension
        ext = os.path.splitext(file_path)[1].lower()
        mime_map = {
            '.pdf': 'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.doc': 'application/msword',
            '.txt': 'text/plain'
        }
        
        return mime_map.get(ext, 'application/octet-stream')
        
    except Exception as e:
        logger.error(f"Error detecting file type: {e}")
        # Final fallback to extension
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pdf':
            return 'application/pdf'
        elif ext in ['.docx', '.doc']:
            return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        elif ext == '.txt':
            return 'text/plain'
        else:
            return 'application/octet-stream'

def read_document_content(filepath: str) -> str:
    """Enhanced document text extraction with automatic file type detection."""
    try:
        if not os.path.exists(filepath):
            raise FileNotFoundError(f"File not found: {filepath}")
        
        # Detect file type
        mime_type = detect_file_type(filepath)
        logger.info(f"Detected file type: {mime_type} for file: {filepath}")
        
        # Extract text based on file type
        if mime_type == "application/pdf":
            return extract_text_from_pdf(filepath)
        elif mime_type in [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword"
        ]:
            return extract_text_from_docx(filepath)
        elif mime_type == "text/plain":
            return extract_text_from_txt(filepath)
        else:
            # Try to determine by file extension as fallback
            ext = os.path.splitext(filepath)[1].lower()
            if ext == '.pdf':
                return extract_text_from_pdf(filepath)
            elif ext in ['.docx', '.doc']:
                return extract_text_from_docx(filepath)
            elif ext == '.txt':
                return extract_text_from_txt(filepath)
            else:
                raise ValueError(f"Unsupported file type: {mime_type}. Supported types: PDF, DOCX, TXT")
            
    except Exception as e:
        logger.error(f"Error reading document {filepath}: {e}")
        raise Exception(f"Error reading document: {str(e)}")

def generate_presentation_content(config: Dict) -> str:
    try:
        print(f"DEBUG: Config keys: {list(config.keys())}")
        print(f"DEBUG: Has doc_content: {'doc_content' in config}")
        if config.get('doc_content'):
            print(f"DEBUG: Document content length: {len(config['doc_content'])}")
        else:
            print(f"DEBUG: Topic: '{config.get('topic', 'NONE')}'")
        
        llm = ChatGoogleGenerativeAI(
            model="gemini-2.0-flash-exp",
            temperature=0.0
        )
        
        if config.get('doc_content'):
            logger.info(f"Generating content from document using few-shot examples")
            
            few_shot_prompt = create_enhanced_few_shot_prompt()
            
            from langchain.prompts import PromptTemplate
            
            doc_template = f"""Here are examples of how to create presentations:

{few_shot_prompt.prefix}

{chr(10).join([few_shot_prompt.example_prompt.format(**example) for example in few_shot_prompt.examples])}

Now, create a {{num_slides}}-slide PowerPoint presentation by summarizing the following document content.
Follow the exact same format as shown in the examples above.

Document Content:
{{doc_content}}

Topic: Document Summary
Number of slides: {{num_slides}}
Tone: {{tone}}
Target audience: {{audience}}
Theme: {{theme}}

Create a {{num_slides}}-slide presentation following the exact format shown. Include image_query and flowchart_description where appropriate.
"""
            
            doc_prompt = PromptTemplate(
                input_variables=["doc_content", "num_slides", "tone", "audience", "theme"],
                template=doc_template
            )
            
            chain = doc_prompt | llm
            output = chain.invoke({
                "doc_content": config["doc_content"],
                "num_slides": config["num_slides"],
                "tone": config["tone"],
                "audience": config["audience"],
                "theme": config["theme"]
            })
            
        else:
            logger.info(f"Generating content for topic: {config.get('topic')}")
            
            few_shot_prompt = create_enhanced_few_shot_prompt()
            chain = few_shot_prompt | llm
            
            prompt_input = {
                "topic": config["topic"],
                "num_slides": config["num_slides"],
                "tone": config["tone"],
                "audience": config["audience"],
                "theme": config["theme"]
            }
            logger.info(f"Using few-shot prompt with input: {prompt_input}")
            
            output = chain.invoke(prompt_input)
        
        logger.info(f"LLM output type: {type(output)}")
        logger.info(f"LLM output: {str(output)[:500]}...")
        
        extracted_text = extract_text_from_llm_output(output)
        logger.info(f"Extracted text length: {len(extracted_text) if extracted_text else 0}")
        
        return extracted_text
        
    except Exception as e:
        logger.error(f"Error generating content: {e}")
        logger.error(f"Error traceback: {traceback.format_exc()}")
        raise

def process_presentation_job(job_id: str, config: Dict):
    try:
        print(f"Starting job processing for job_id: {job_id}")
        
        if job_id not in active_jobs:
            print(f"WARNING: Job {job_id} not found in active_jobs at start of processing")
            active_jobs[job_id] = {
                'status': 'started',
                'progress': 0,
                'created_at': datetime.now().isoformat(),
                'config': config
            }
        
        active_jobs[job_id]['status'] = 'generating_content'
        active_jobs[job_id]['progress'] = 20
        print(f"Job {job_id}: Updated status to generating_content")
        
        logger.info(f"Job {job_id}: Generating AI content")
        content = generate_presentation_content(config)
        
        active_jobs[job_id]['status'] = 'parsing_content'
        active_jobs[job_id]['progress'] = 40
        
        logger.info(f"Job {job_id}: Parsing slide content")
        slides_data = parse_slides_enhanced(content)
        
        if not slides_data:
            raise Exception("Could not parse slide content")
        
        active_jobs[job_id]['status'] = 'creating_presentation'
        active_jobs[job_id]['progress'] = 60
        
        logger.info(f"Job {job_id}: Creating presentation")
        
        if os.environ.get('DYNO'):
            output_folder = '/tmp'
        else:
            output_folder = 'outputs'
            os.makedirs(output_folder, exist_ok=True)
        
        file_format = config.get('file_format', 'pptx')
        
        user_filename = config['filename']
        
        base_filename = user_filename.split('.')[0] if '.' in user_filename else user_filename
        
        if file_format == 'pdf':
            final_extension = '.pdf'
        else:
            final_extension = '.pptx'
            file_format = 'pptx'
        
        if os.environ.get('DYNO'):
            output_filename = f"/tmp/{job_id}_{base_filename}"
        else:
            output_filename = os.path.join(output_folder, f"{job_id}_{base_filename}")
        
        processed_flowcharts = []
        flowcharts = config.get('flowcharts', [])
        logger.info(f"Processing {len(flowcharts)} flowcharts")
        
        for i, flowchart in enumerate(flowcharts):
            logger.info(f"Processing flowchart {i+1}: {flowchart}")
            if isinstance(flowchart, dict) and 'description' in flowchart:
                original_desc = flowchart['description']
                logger.info(f"Original flowchart description: {original_desc[:100]}...")
                
                converted_desc = convert_simple_flowchart_to_mermaid(original_desc)
                if converted_desc:
                    processed_flowcharts.append({
                        'title': flowchart.get('title', 'Flowchart'),
                        'description': converted_desc
                    })
                    logger.info(f"Successfully converted simple flowchart format to Mermaid")
                    logger.info(f"Converted description: {converted_desc[:100]}...")
                else:
                    processed_flowcharts.append(flowchart)
                    logger.warning(f"Flowchart conversion failed, using original format")
            else:
                processed_flowcharts.append(flowchart)
                logger.info(f"Using flowchart as-is (not dict or missing description)")
        
        logger.info(f"Final processed flowcharts: {len(processed_flowcharts)}")
        
        flowcharts_tuples = []
        for flowchart in processed_flowcharts:
            if isinstance(flowchart, dict):
                title = flowchart.get('title', 'Flowchart')
                description = flowchart.get('description', '')
                flowcharts_tuples.append((title, description))
                logger.info(f"Converted flowchart: {title}")
            else:
                logger.warning(f"Unexpected flowchart format: {type(flowchart)}")
        
        logger.info(f"Final flowcharts for PPT creation: {len(flowcharts_tuples)} tuples")
        
        create_enhanced_ppt(
            slides_data=slides_data,
            filename=output_filename,
            output_format=file_format,
            theme=config['theme'],
            text_size=config['text_size'],
            flowcharts=flowcharts_tuples
        )
        
        if config.get('file_format') == 'pdf':
            logger.info(f"Job {job_id}: Converting to PDF")
            active_jobs[job_id]['status'] = 'converting_to_pdf'
            active_jobs[job_id]['progress'] = 80
            
            pptx_file = f"{output_filename}.pptx"
            pdf_file = f"{output_filename}.pdf"
            
            logger.info(f"PPTX file path: {pptx_file}")
            logger.info(f"PDF file path: {pdf_file}")
            logger.info(f"PPTX exists: {os.path.exists(pptx_file)}")
            
            try:
                if not os.path.exists(pptx_file):
                    raise Exception(f"PPTX file not found for conversion: {pptx_file}")
                
                logger.info(f"Converting PPTX to PDF: {pptx_file} -> {pdf_file}")
                
                success = convert_pptx_to_pdf(pptx_file, pdf_file)
                
                if not success:
                    raise Exception("PDF conversion returned failure status")
                
                if not os.path.exists(pdf_file):
                    raise Exception("PDF file was not created")
                
                pdf_size = os.path.getsize(pdf_file)
                if pdf_size == 0:
                    raise Exception("PDF file is empty")
                
                try:
                    with open(pdf_file, 'rb') as f:
                        header = f.read(10)
                        if not header.startswith(b'%PDF-'):
                            raise Exception(f"Invalid PDF file format. Header: {header}")
                except Exception as verify_error:
                    raise Exception(f"PDF verification failed: {verify_error}")
                
                logger.info(f"PDF created successfully: {pdf_file} ({pdf_size} bytes)")
                
                try:
                    if os.path.exists(pptx_file):
                        os.remove(pptx_file)
                        logger.info(f"Removed temporary PPTX file: {pptx_file}")
                except Exception as cleanup_error:
                    logger.warning(f"Could not remove temporary PPTX file: {cleanup_error}")
                    
                final_output_filename = f"{job_id}_{base_filename}.pdf"
                final_filepath = pdf_file
                logger.info(f"PDF conversion successful: {final_filepath}")
                
            except Exception as pdf_error:
                logger.error(f"PDF conversion failed: {pdf_error}")
                import traceback
                logger.error(f"PDF conversion traceback: {traceback.format_exc()}")
                logger.warning("Falling back to PPTX format")
                
                try:
                    if os.path.exists(pdf_file):
                        os.remove(pdf_file)
                except:
                    pass
                
                final_output_filename = f"{job_id}_{base_filename}.pptx"
                final_filepath = pptx_file
        else:
            final_output_filename = f"{job_id}_{base_filename}.pptx"
            if os.environ.get('DYNO'):
                final_filepath = f"/tmp/{job_id}_{base_filename}.pptx"
            else:
                final_filepath = f"{output_filename}.pptx"
        
        active_jobs[job_id]['status'] = 'completed'
        active_jobs[job_id]['progress'] = 100
        
        job_results[job_id] = {
            'status': 'completed',
            'filename': final_output_filename,
            'filepath': final_filepath,
            'slides_count': len(slides_data),
            'created_at': datetime.now().isoformat(),
            'config': config
        }
        
        logger.info(f"Job {job_id}: Completed successfully")
        print(f"Job {job_id}: Moved from active_jobs to job_results")
        print(f"Active jobs remaining: {list(active_jobs.keys())}")
        print(f"Completed jobs: {list(job_results.keys())}")
        
        if job_id in active_jobs:
            del active_jobs[job_id]
        
    except Exception as e:
        logger.error(f"Job {job_id} failed: {e}")
        print(f"Job {job_id} FAILED: {e}")
        
        if job_id in active_jobs:
            active_jobs[job_id]['status'] = 'failed'
            active_jobs[job_id]['error'] = str(e)
        
        job_results[job_id] = {
            'status': 'failed',
            'error': str(e),
            'created_at': datetime.now().isoformat()
        }

class PPTGeneratorAPI:
    
    def __init__(self, base_url: str = "http://localhost:5000"):
        self.base_url = base_url.rstrip('/')
    
    def health_check(self):
        import requests
        try:
            response = requests.get(f"{self.base_url}/api/health")
            return response.status_code == 200
        except:
            return False
    
    def get_config(self):
        import requests
        response = requests.get(f"{self.base_url}/api/config")
        return response.json() if response.status_code == 200 else None
    
    def upload_document(self, file_path: str):
        import requests
        with open(file_path, 'rb') as f:
            files = {'file': f}
            response = requests.post(f"{self.base_url}/api/upload", files=files)
        return response.json()
    
    def generate_presentation(self, config: Dict):
        import requests
        response = requests.post(f"{self.base_url}/api/generate", json=config)
        return response.json()
    
    def get_job_status(self, job_id: str):
        import requests
        response = requests.get(f"{self.base_url}/api/status/{job_id}")
        return response.json()
    
    def download_presentation(self, job_id: str, save_path: str):
        import requests
        response = requests.get(f"{self.base_url}/api/download/{job_id}")
        if response.status_code == 200:
            with open(save_path, 'wb') as f:
                f.write(response.content)
            return True
        return False

def convert_simple_flowchart_to_mermaid(description: str) -> str:
    try:
        logger.info(f"Converting flowchart: {description[:50]}...")
        
        if description.strip().lower().startswith(('flowchart', 'graph', 'sequencediagram')):
            logger.info("Input appears to be valid Mermaid syntax, returning as-is")
            return description
        
        description = description.strip()
        if not description:
            logger.warning("Empty flowchart description provided")
            return None
        
        connections = []
        
        parts = re.split(r'[;\n,]', description)
        logger.info(f"Split into {len(parts)} parts: {parts}")
        
        nodes = set()
        edges = []
        
        for part in parts:
            part = part.strip()
            if not part:
                continue
                
            if '->' in part:
                arrow_parts = part.split('->')
            elif '→' in part:
                arrow_parts = part.split('→')
            elif '-->' in part:
                arrow_parts = part.split('-->')
            elif '=>' in part:
                arrow_parts = part.split('=>')
            else:
                logger.warning(f"No arrows found in part: {part}")
                continue
                
            for i in range(len(arrow_parts) - 1):
                from_node = arrow_parts[i].strip()
                to_node = arrow_parts[i + 1].strip()
                
                if from_node and to_node:
                    from_clean = re.sub(r'[^\w\s-]', '', from_node).strip()
                    to_clean = re.sub(r'[^\w\s-]', '', to_node).strip()
                    
                    from_id = re.sub(r'\s+', '_', from_clean)
                    to_id = re.sub(r'\s+', '_', to_clean)
                    
                    if from_id and to_id:
                        nodes.add((from_id, from_clean))
                        nodes.add((to_id, to_clean))
                        edges.append((from_id, to_id))
                        logger.debug(f"Added edge: {from_id} -> {to_id}")
        
        if not edges:
            logger.warning("No valid edges found in flowchart description")
            return None
        
        logger.info(f"Found {len(nodes)} nodes and {len(edges)} edges")
        
        mermaid_lines = ['flowchart TD']
        
        for node_id, node_label in nodes:
            if node_label.lower() in ['start', 'begin']:
                mermaid_lines.append(f'    {node_id}([{node_label}])')
            elif node_label.lower() in ['end', 'finish', 'stop']:
                mermaid_lines.append(f'    {node_id}([{node_label}])')
            elif '?' in node_label or 'decision' in node_label.lower():
                mermaid_lines.append(f'    {node_id}{{{node_label}}}')
            else:
                mermaid_lines.append(f'    {node_id}[{node_label}]')
        
        for from_id, to_id in edges:
            mermaid_lines.append(f'    {from_id} --> {to_id}')
        
        result = '\n'.join(mermaid_lines)
        logger.info(f"Successfully converted to Mermaid: {len(mermaid_lines)} lines")
        return result
        
    except Exception as e:
        logger.error(f"Error converting simple flowchart: {e}")
        logger.error(f"Error traceback: {traceback.format_exc()}")
        return None