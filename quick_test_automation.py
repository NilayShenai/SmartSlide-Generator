#!/usr/bin/env python3
"""
Quick test of LibreOffice automation in application context
"""

import os
import sys
from pathlib import Path

# Add the app directory to the Python path
current_dir = Path(__file__).parent
app_dir = current_dir / "app"
sys.path.insert(0, str(app_dir))

try:
    from utils import convert_pptx_to_pdf
    from pptx import Presentation
    
    print("🧪 Quick LibreOffice Automation Test")
    print("=" * 40)
    
    # Create a simple test presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Automation Test"
    prs.save('quick_test.pptx')
    print("✅ Created test PPTX")
    
    # Convert using the automated LibreOffice function
    print("🔄 Converting with automated LibreOffice...")
    result = convert_pptx_to_pdf('quick_test.pptx', 'quick_test.pdf')
    
    # Check results
    if result and os.path.exists('quick_test.pdf'):
        file_size = os.path.getsize('quick_test.pdf')
        print(f"✅ Conversion successful!")
        print(f"📄 PDF size: {file_size:,} bytes")
        print("🎉 LibreOffice automation working perfectly!")
    else:
        print("❌ Conversion failed")
        
    # Cleanup
    for file in ['quick_test.pptx', 'quick_test.pdf']:
        if os.path.exists(file):
            os.remove(file)
            print(f"🗑️ Cleaned up {file}")
            
except ImportError as e:
    print(f"❌ Import error: {e}")
except Exception as e:
    print(f"❌ Error: {e}")
    import traceback
    traceback.print_exc()
